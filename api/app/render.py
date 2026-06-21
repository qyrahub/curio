"""Export rendering. Async job pattern so a 50-page render never blocks a
request: POST /plans/{id}/exports queues a job, the client polls
GET /exports/{id}, then downloads when status == done.

PPTX is a real renderer (python-pptx). PDF / DOCX / JPEG are the next increment;
they currently write a placeholder so the pipeline is exercised end to end.
A PDF/JPEG path can be derived from the PPTX via LibreOffice on the VPS later.
"""
from __future__ import annotations

import json
import os

from .config import settings
from .models import Plan
from .store import get_export, save_export
from . import images

EXPORT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "exports")
os.makedirs(EXPORT_DIR, exist_ok=True)

# ---- palette (Curio brand) ----
INK = "2C2A4A"
WHITE = "FFFFFF"
PAPER = "FFFBF4"
SUN = "FFC94D"
GOLD = "9A6B00"
MUTED = "6A6680"

# Per-subject colours, themed by child gender (neutral = the varied default).
DOMAIN_NEUTRAL = {
    "biology": "5BBF8A", "maths": "5AA7E6", "language": "9B6DD6",
    "geography": "2EC4B6", "history": "C0894A", "health": "FF7A66",
    "technology": "6C7AE0", "arts": "E0699B", "science": "3FA79E",
    "fun": "FFC94D", "speech": "FF7A66",
}
DOMAIN_GIRL = {
    "biology": "E0699B", "maths": "9B6DD6", "language": "C45CB0",
    "geography": "B673D6", "history": "D17AA0", "health": "FF7AA8",
    "technology": "8E6AD8", "arts": "E0699B", "science": "C779B0",
    "fun": "FFB3D1", "speech": "FF7AA8",
}
DOMAIN_BOY = {
    "biology": "5BBF8A", "maths": "5AA7E6", "language": "4C7AE0",
    "geography": "2EC4B6", "history": "3F8FD0", "health": "5AA7E6",
    "technology": "4C7AE0", "arts": "3FA79E", "science": "2EC4B6",
    "fun": "5AA7E6", "speech": "5AA7E6",
}
ACCENT = {"neutral": "FF7A66", "girl": "C45CB0", "boy": "4C7AE0"}
# A festive confetti mix per theme, used on the cover.
CONFETTI = {
    "neutral": ["FF7A66", "FFC94D", "5BBF8A", "5AA7E6", "9B6DD6", "2EC4B6"],
    "girl": ["FF7AA8", "C45CB0", "9B6DD6", "FFB3D1", "E0699B", "FFC94D"],
    "boy": ["4C7AE0", "5AA7E6", "2EC4B6", "5BBF8A", "6C7AE0", "FFC94D"],
}


def _gender(req) -> str:
    g = (getattr(req, "gender", None) or "neutral").lower()
    return g if g in ("girl", "boy", "neutral") else "neutral"


def _palette(gender: str) -> dict:
    return {"girl": DOMAIN_GIRL, "boy": DOMAIN_BOY}.get(gender, DOMAIN_NEUTRAL)


def _label(domain: str) -> str:
    if domain == "fun":
        return "FUN BREAK"
    return domain.upper()


def export_path(plan_id: str, fmt: str) -> str:
    return os.path.join(EXPORT_DIR, f"{plan_id}.{fmt}")


def download_name(plan_id: str, fmt: str) -> str:
    return f"curio-plan-{plan_id}.{fmt}"


def plan_meta(plan: Plan) -> dict:
    pages = [{"title": p.topic, "text": p.guideline, "challenge": p.challenge, "answer": getattr(p, "answer", "") or ""}
             for p in plan.pages if getattr(p, "archetype", "") != "fun"]
    full = [f"Curio learning plan — age {plan.request.age}", ""]
    for i, pg in enumerate(pages, 1):
        full += [f"{i}. {pg['title']}", pg["text"], f"Try this: {pg['challenge']}", ""]
    return {
        "title": f"Curio plan — age {plan.request.age}",
        "subtitle": f"{len(pages)}-page learning adventure",
        "blurb": "A bright, balanced, age-matched learning plan from Curio.",
        "pages": pages, "full_text": "\n".join(full), "slug": f"plan-{plan.id}", "emoji": "✨", "color": "5AA7E6",
    }


# ---------------- PPTX renderer ----------------
def _render_pptx(plan: Plan, path: str) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    FONT = "Calibri"

    def rgb(h):
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def tint(h, f):
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return RGBColor(int(r + (255 - r) * f), int(g + (255 - g) * f), int(b + (255 - b) * f))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def bg(slide, hexc):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(hexc)

    def round_rect(slide, x, y, w, h, fill_hex, radius=0.12):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill_hex) if isinstance(fill_hex, str) else fill_hex
        shp.line.fill.background()
        shp.shadow.inherit = False
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
        return shp

    def circle(slide, x, y, d, fill_hex):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill_hex)
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(2)
        tf.margin_top = tf.margin_bottom = Pt(2)
        tf.vertical_anchor = anchor
        first = True
        for (s, size, bold, color) in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            r = p.add_run()
            r.text = s
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = rgb(color)
        return tb

    def shape_text(shp, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        first = True
        for (s, size, bold, color) in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            r = p.add_run()
            r.text = s
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = rgb(color)

    # ---- cover slide (light, vibrant, colourful) ----
    req = plan.request
    gender = _gender(req)
    pal = _palette(gender)
    accent = ACCENT[gender]
    confetti = CONFETTI[gender]

    s = prs.slides.add_slide(blank)
    bg(s, PAPER)
    # scattered confetti around the edges (kept clear of the centred text band)
    spots = [
        (0.7, 0.7, 1.1), (11.4, 0.6, 0.8), (12.2, 1.7, 0.5), (0.5, 2.2, 0.6),
        (11.9, 5.2, 0.9), (0.8, 6.1, 0.7), (2.0, 6.6, 0.45), (10.6, 6.6, 0.6),
        (12.5, 3.4, 0.4), (0.4, 4.4, 0.45),
    ]
    for i, (x, y, d) in enumerate(spots):
        circle(s, x, y, d, confetti[i % len(confetti)])
    # star mark
    circle(s, 6.05, 1.5, 1.0, SUN)
    text(s, 6.05, 1.52, 1.0, 1.0, [("★", 34, True, accent)], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    age = "12+" if req.age >= 13 else str(req.age)
    headline = "A speech plan" if req.cadence == "speech" else "A learning adventure"
    text(s, 0.9, 2.95, 11.5, 1.6, [("Curio", 56, True, INK)], PP_ALIGN.CENTER)
    text(s, 0.95, 4.25, 11.5, 1.0, [(f"{headline} for age {age}", 26, True, accent)], PP_ALIGN.CENTER)
    bits = [req.cadence]
    if req.subjects:
        bits.append("subject: " + req.subjects[0])
    if req.interests:
        bits.append(", ".join([i for i in req.interests if i != "surprise me"][:3]) or "surprise mix")
    text(s, 0.95, 5.05, 11.5, 0.8, [("  ·  ".join([b for b in bits if b]), 16, False, MUTED)], PP_ALIGN.CENTER)

    # ---- content slides ----
    total = len(plan.pages)
    for pg in plan.pages:
        s = prs.slides.add_slide(blank)
        bg(s, WHITE)
        dcol = pal.get(pg.domain, "5AA7E6")
        emoji, _, desc = pg.anchor_visual.partition("|")
        text(s, 0.6, 0.42, 6.0, 0.4, [("Curio", 12, True, MUTED)])
        text(s, 11.0, 6.95, 1.7, 0.4, [(f"{pg.order} of {total}", 11, True, MUTED)], PP_ALIGN.RIGHT)

        show_visual = req.medium != "text"
        show_text = req.medium != "illustration"
        right_x = 6.9 if show_visual else 0.6
        right_w = 5.8 if show_visual else 12.1

        if show_visual:
            round_rect(s, 0.6, 1.4, 5.7, 5.3, tint(dcol, 0.86), radius=0.10)
            img_path = images.image_for_page(pg.topic, desc, age=req.age)
            if img_path:
                # real generated illustration, framed by the tinted panel
                s.shapes.add_picture(img_path, Inches(1.6), Inches(1.7), Inches(3.7), Inches(3.7))
            else:
                # graceful fallback to the emoji motif when generation is off
                circle(s, 2.25, 2.15, 2.4, dcol)
                text(s, 2.25, 2.15, 2.4, 2.4, [(emoji or "✨", 60, False, WHITE)], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            if desc:
                text(s, 0.95, 5.55, 5.0, 1.0, [(desc, 13, False, MUTED)], PP_ALIGN.CENTER)

        badge = round_rect(s, right_x, 1.4, 2.8, 0.45, dcol, radius=0.5)
        shape_text(badge, [(_label(pg.domain), 12, True, WHITE)], PP_ALIGN.CENTER)
        text(s, right_x, 2.05, right_w, 1.3, [(pg.topic, 34, True, INK)])
        if show_text:
            text(s, right_x, 3.35, right_w, 1.5, [(pg.guideline, 16, False, INK)])
        ch = round_rect(s, right_x, 5.05, right_w, 1.5, tint(SUN, 0.80), radius=0.10)
        shape_text(ch, [("TRY THIS", 11, True, GOLD), (pg.challenge, 15, False, INK)], PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

    # ---- closing slide (light) ----
    s = prs.slides.add_slide(blank)
    bg(s, PAPER)
    for i, (x, y, d) in enumerate([(0.8, 0.8, 0.9), (11.6, 0.7, 0.7), (12.0, 5.6, 0.8), (0.7, 6.0, 0.6)]):
        circle(s, x, y, d, confetti[i % len(confetti)])
    text(s, 0.9, 3.0, 11.5, 1.4, [("Made with Curio", 40, True, INK)], PP_ALIGN.CENTER)
    text(s, 0.9, 4.2, 11.5, 0.8, [("Keep exploring, together. \u2728", 20, True, accent)], PP_ALIGN.CENTER)

    prs.save(path)


def _render_placeholder(plan: Plan, fmt: str, path: str) -> None:
    data = {
        "note": f"PLACEHOLDER {fmt} export. Real {fmt} renderer not wired yet.",
        "plan_id": plan.id,
        "pages": [p.model_dump() for p in plan.pages],
    }
    with open(path, "w", encoding="utf-8") as fh:
        json.dump(data, fh, indent=2)


def run_export(export_id: str, plan: Plan) -> None:
    job = get_export(export_id)
    if job is None:
        return
    job.status = "running"
    save_export(job)
    try:
        path = export_path(plan.id, job.fmt)
        if job.fmt == "pptx":
            _render_pptx(plan, path)
        elif job.fmt in ("pdf", "epub", "txt", "md"):
            from . import bookexport
            with open(path, "wb") as fh:
                fh.write(bookexport.render_doc(plan_meta(plan), job.fmt))
        else:
            _render_placeholder(plan, job.fmt, path)
        job.download_url = f"{settings.public_base}/exports/{export_id}/download"
        job.status = "done"
    except Exception as exc:  # noqa: BLE001
        job.status = "error"
        job.error = str(exc)
    save_export(job)
