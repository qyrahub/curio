"""Generate real downloadable files (PDF, PPTX, EPUB, TXT, MD) from a generic
document 'meta' — used for both Library books and Child plans.
"""
from __future__ import annotations

import html
import io
import zipfile

from .models import BookDetail

MEDIA = {
    "pdf": "application/pdf",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "epub": "application/epub+zip",
    "txt": "text/plain; charset=utf-8",
    "md": "text/markdown; charset=utf-8",
}


def _slugify(s: str) -> str:
    return "".join(ch if ch.isalnum() else "-" for ch in (s or "").lower()).strip("-") or "curio"


def book_meta(d: BookDetail) -> dict:
    return {
        "title": d.title,
        "subtitle": f"{d.category_label} · {d.discipline} · ages {d.age_min}-{d.age_max}",
        "blurb": d.blurb,
        "pages": [{"title": p.title, "text": p.text, "challenge": p.challenge, "answer": p.answer} for p in d.book_pages],
        "full_text": d.full_text,
        "slug": _slugify(d.title),
        "emoji": d.emoji,
        "color": d.color,
    }


def filename(meta: dict, fmt: str) -> str:
    return f"{meta.get('slug') or 'curio'}.{fmt}"


def _md(m: dict) -> str:
    out = [f"# {m['title']}", "", f"*{m.get('blurb','')}*", "", f"**{m.get('subtitle','')}**", "", "---", ""]
    for i, p in enumerate(m["pages"], 1):
        out += [f"## {i}. {p['title']}", "", p["text"], "", f"> **Try this:** {p['challenge']}", ""]
        if p.get("answer"):
            out += [f"> _Answer: {p['answer']}_", ""]
    out += ["---", "", "Made with Curio."]
    return "\n".join(out)


def _txt(m: dict) -> str:
    return m.get("full_text") or _md(m).replace("#", "").replace("*", "").replace(">", "")


def _pdf(m: dict) -> bytes:
    from reportlab.lib.pagesizes import A5
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.lib.enums import TA_CENTER
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A5, topMargin=18 * mm, bottomMargin=18 * mm,
                            leftMargin=16 * mm, rightMargin=16 * mm, title=m["title"])
    ss = getSampleStyleSheet()
    title = ParagraphStyle("t", parent=ss["Title"], fontSize=24, leading=28, textColor="#2C2A4A", alignment=TA_CENTER)
    sub = ParagraphStyle("s", parent=ss["Normal"], fontSize=11, leading=15, textColor="#6b6b80", alignment=TA_CENTER, spaceBefore=8)
    h2 = ParagraphStyle("h2", parent=ss["Heading2"], fontSize=15, leading=19, textColor="#F2563D", spaceBefore=6, spaceAfter=8)
    body = ParagraphStyle("b", parent=ss["Normal"], fontSize=12, leading=18, textColor="#2C2A4A")
    chal = ParagraphStyle("c", parent=ss["Normal"], fontSize=11, leading=16, textColor="#2EA37A", spaceBefore=8, backColor="#EAF7F0", borderPadding=6)

    def esc(s: str) -> str:
        return html.escape(s or "")

    story = [Spacer(1, 38 * mm), Paragraph(esc(m["title"]), title)]
    if m.get("blurb"):
        story.append(Paragraph(esc(m["blurb"]), sub))
    if m.get("subtitle"):
        story.append(Paragraph(esc(m["subtitle"]), sub))
    story.append(PageBreak())
    for i, p in enumerate(m["pages"], 1):
        story += [Paragraph(f"{i}. {esc(p['title'])}", h2), Paragraph(esc(p["text"]), body),
                  Paragraph(f"<b>Try this:</b> {esc(p['challenge'])}", chal)]
        if p.get("answer"):
            story.append(Paragraph(f"<i>Answer: {esc(p['answer'])}</i>", body))
        story.append(PageBreak())
    story.append(Paragraph("Made with Curio.", sub))
    doc.build(story)
    return buf.getvalue()


def _pptx(m: dict) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    color = (m.get("color") or "5AA7E6").lstrip("#")
    emoji = m.get("emoji") or "✨"
    INK = "2C2A4A"
    WHITE = RGBColor(255, 255, 255)

    def rgb(h):
        h = h.lstrip("#"); return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def tint(h, f):
        h = h.lstrip("#"); r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return RGBColor(int(r + (255 - r) * f), int(g + (255 - g) * f), int(b + (255 - b) * f))

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def bg(s, c):
        s.background.fill.solid(); s.background.fill.fore_color.rgb = c

    def rrect(s, x, y, w, h, fill, rad=0.12):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background(); shp.shadow.inherit = False
        try: shp.adjustments[0] = rad
        except Exception: pass
        return shp

    def circle(s, x, y, d, fill):
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background(); shp.shadow.inherit = False
        return shp

    def tx(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        first = True
        for (t, size, bold, col) in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.alignment = align
            r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col
        return tb

    # ---- title slide ----
    s = prs.slides.add_slide(blank); bg(s, tint(color, 0.84))
    circle(s, 5.67, 1.25, 2.0, rgb(color))
    tx(s, 4.67, 1.3, 4.0, 1.9, [(emoji, 90, False, WHITE)], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    tx(s, 1.2, 3.7, 10.93, 1.6, [(m["title"], 40, True, rgb(INK))], PP_ALIGN.CENTER)
    if m.get("blurb"):
        tx(s, 1.5, 5.25, 10.33, 1.0, [(m["blurb"], 18, False, rgb("6B6B80"))], PP_ALIGN.CENTER)
    if m.get("subtitle"):
        tx(s, 1.5, 6.1, 10.33, 0.6, [(m["subtitle"], 14, False, rgb("9A93A8"))], PP_ALIGN.CENTER)

    # ---- chapter slides ----
    for i, p in enumerate(m["pages"], 1):
        s = prs.slides.add_slide(blank); bg(s, rgb("FFFBF4"))
        rrect(s, -0.3, -0.3, 4.6, 8.1, tint(color, 0.72), 0.0)   # left panel
        circle(s, 1.4, 1.2, 1.6, rgb(color))
        tx(s, 0.9, 1.25, 2.6, 1.5, [(emoji, 60, False, WHITE)], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        tx(s, 0.4, 3.1, 3.5, 2.0, [(f"Chapter {i}", 18, True, rgb(color)), (p["title"], 24, True, rgb(INK))], PP_ALIGN.CENTER)
        tx(s, 4.9, 0.7, 7.9, 1.3, [(f"{i}. {p['title']}", 30, True, rgb(INK))])
        tx(s, 4.9, 2.0, 7.9, 3.0, [(p["text"], 18, False, rgb(INK))])
        rrect(s, 4.9, 5.25, 7.9, 1.6, tint("2EC4B6", 0.82), 0.18)
        tx(s, 5.2, 5.4, 7.3, 1.3, [("Try this", 13, True, rgb("1F8E80")), (p["challenge"], 16, False, rgb(INK))], PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

    buf = io.BytesIO(); prs.save(buf)
    return buf.getvalue()


def _epub(m: dict) -> bytes:
    buf = io.BytesIO()
    uid = f"curio-{m.get('slug', 'book')}"

    def xhtml(title: str, inner: str) -> str:
        return (f'<?xml version="1.0" encoding="utf-8"?>\n<!DOCTYPE html>'
                f'<html xmlns="http://www.w3.org/1999/xhtml"><head><title>{html.escape(title)}</title></head>'
                f'<body>{inner}</body></html>')

    chapters = [("cover", "Cover", xhtml(m["title"],
                f"<h1>{html.escape(m['title'])}</h1><p><em>{html.escape(m.get('blurb',''))}</em></p>"
                f"<p>{html.escape(m.get('subtitle',''))}</p>"))]
    for i, p in enumerate(m["pages"], 1):
        inner = (f"<h2>{i}. {html.escape(p['title'])}</h2><p>{html.escape(p['text'])}</p>"
                 f"<p><strong>Try this:</strong> {html.escape(p['challenge'])}</p>")
        chapters.append((f"chap{i}", f"{i}. {p['title']}", xhtml(p["title"], inner)))

    manifest = "\n".join(f'<item id="{c}" href="{c}.xhtml" media-type="application/xhtml+xml"/>' for c, _t, _x in chapters)
    spine = "\n".join(f'<itemref idref="{c}"/>' for c, _t, _x in chapters)
    opf = (f'<?xml version="1.0" encoding="utf-8"?>\n<package xmlns="http://www.idpf.org/2007/opf" version="2.0" unique-identifier="bookid">'
           f'<metadata xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:identifier id="bookid">{uid}</dc:identifier>'
           f'<dc:title>{html.escape(m["title"])}</dc:title><dc:language>en</dc:language><dc:creator>Curio</dc:creator></metadata>'
           f'<manifest>{manifest}<item id="ncx" href="toc.ncx" media-type="application/x-dtbncx+xml"/></manifest>'
           f'<spine toc="ncx">{spine}</spine></package>')
    navpoints = "\n".join(
        f'<navPoint id="{c}" playOrder="{i+1}"><navLabel><text>{html.escape(t)}</text></navLabel><content src="{c}.xhtml"/></navPoint>'
        for i, (c, t, _x) in enumerate(chapters))
    ncx = (f'<?xml version="1.0" encoding="utf-8"?>\n<ncx xmlns="http://www.daisy.org/z3986/2005/ncx/" version="2005-1">'
           f'<head><meta name="dtb:uid" content="{uid}"/></head><docTitle><text>{html.escape(m["title"])}</text></docTitle>'
           f'<navMap>{navpoints}</navMap></ncx>')
    container = ('<?xml version="1.0"?>\n<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
                 '<rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/></rootfiles></container>')
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("mimetype", "application/epub+zip", compress_type=zipfile.ZIP_STORED)
        z.writestr("META-INF/container.xml", container)
        z.writestr("OEBPS/content.opf", opf)
        z.writestr("OEBPS/toc.ncx", ncx)
        for c, _t, x in chapters:
            z.writestr(f"OEBPS/{c}.xhtml", x)
    return buf.getvalue()


def render_doc(meta: dict, fmt: str) -> bytes:
    if fmt == "pdf":
        return _pdf(meta)
    if fmt == "pptx":
        return _pptx(meta)
    if fmt == "epub":
        return _epub(meta)
    if fmt == "md":
        return _md(meta).encode("utf-8")
    return _txt(meta).encode("utf-8")


def render(d: BookDetail, fmt: str) -> bytes:
    """Backwards-compatible Library entrypoint."""
    return render_doc(book_meta(d), fmt)
